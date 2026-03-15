#!/usr/bin/env python3
"""
生成 Neuroscience Daily Digest PowerPoint 演示文稿
使用简约学术风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

class PPTGenerator:
    """PPT生成器"""
    
    COLORS = {
        'primary': (26, 54, 93),
        'secondary': (44, 82, 130),
        'text': (45, 55, 72),
        'muted': (74, 85, 104),
        'light': (113, 128, 150),
    }
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def _add_textbox(self, slide, left, top, width, height, text, 
                     font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                     font_name='Microsoft YaHei'):
        """添加文本框"""
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = RGBColor(*color)
        
        return textbox
    
    def add_title_slide(self, date_str, count):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 图标
        self._add_textbox(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(1),
                         "🧠", font_size=48, align=PP_ALIGN.CENTER)
        
        # 主标题
        self._add_textbox(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1),
                         "Neuroscience", font_size=36, bold=True, 
                         color=self.COLORS['primary'], align=PP_ALIGN.CENTER)
        
        # 副标题
        self._add_textbox(slide, Inches(1), Inches(3.3), Inches(11.333), Inches(0.8),
                         "Daily Digest", font_size=28, 
                         color=self.COLORS['secondary'], align=PP_ALIGN.CENTER)
        
        # 日期
        self._add_textbox(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(0.6),
                         date_str, font_size=18, 
                         color=self.COLORS['light'], align=PP_ALIGN.CENTER)
        
        # 统计
        self._add_textbox(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.6),
                         f"精选 {count} 篇文献", font_size=16, 
                         color=self.COLORS['muted'], align=PP_ALIGN.CENTER)
        
        return slide
    
    def add_content_slide(self, number, en_title, cn_title, abstract, journal, date):
        """添加内容页（单篇文章）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 序号+英文标题
        self._add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8),
                         f"{number}. {en_title}", font_size=20, bold=True,
                         color=self.COLORS['primary'])
        
        # 中文标题
        self._add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5),
                         cn_title, font_size=16, 
                         color=self.COLORS['muted'])
        
        # 摘要
        self._add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4),
                         abstract, font_size=14, 
                         color=self.COLORS['text'])
        
        # 来源
        self._add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
                         f"📄 {journal} | {date}", font_size=12, 
                         color=self.COLORS['light'])
        
        return slide
    
    def save(self, output_path):
        """保存PPT"""
        self.prs.save(output_path)
        print(f"✅ PPT已生成: {output_path}")


def main():
    """主函数"""
    import sys
    
    if len(sys.argv) < 3:
        print("用法: python generate_ppt.py --input input.md --output output.pptx")
        return
    
    input_file = sys.argv[2]
    output_file = sys.argv[4] if len(sys.argv) > 4 else 'neuro_digest.pptx'
    
    # 读取并解析（复用generate_word的解析函数）
    from generate_word import parse_markdown
    
    with open(input_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    articles = parse_markdown(md_content)
    
    # 生成PPT
    gen = PPTGenerator()
    date_str = datetime.now().strftime('%Y年%m月%d日')
    
    # 标题页
    gen.add_title_slide(date_str, len(articles))
    
    # 内容页
    for article in articles:
        gen.add_content_slide(
            number=article.get('number', '1'),
            en_title=article.get('title', ''),
            cn_title=article.get('cn_title', ''),
            abstract=article.get('abstract', '')[:500] + '...',  # PPT摘要精简
            journal=article.get('journal', 'Unknown'),
            date=article.get('date', date_str)
        )
    
    gen.save(output_file)


if __name__ == '__main__':
    main()
